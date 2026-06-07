import io
import json
import os
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from groq import AsyncGroq
from pydantic import BaseModel
from pptx import Presentation
from pypdf import PdfReader

try:
    from . import database
except ImportError:  # python backend/main.py 형태의 직접 실행 호환
    import database

BASE_DIR = Path(__file__).resolve().parent
load_dotenv(dotenv_path=BASE_DIR / ".env")

app = FastAPI()

# Flutter Web 개발용 CORS입니다. 운영 환경에서는 allow_origins를 실제 도메인으로 제한해야 합니다.
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

client = AsyncGroq(api_key=os.getenv("GROQ_API_KEY"))


class QuestionRequest(BaseModel):
    subject: str
    difficulty: str = "beginner"
    level_title: str = ""
    level_description: str = ""
    count: int = 10
    session_id: int = 0
    mode: str = "recommended"
    generation_mode: str = "ai_only"
    learning_level: str = "beginner"
    uses_uploaded_material: bool = False
    subject_type: str = "conceptual"
    level: int = 1
    allowed_question_types: List[str] = []
    question_type_weights: Dict[str, int] = {}


class CurriculumRequest(BaseModel):
    subject: str
    user_id: int = 1
    generation_mode: str = "ai_only"
    difficulty: str = "beginner"
    learning_level: str = "beginner"


class SessionSetupRequest(BaseModel):
    generation_mode: str = "ai_only"
    difficulty: str = "beginner"
    learning_level: str = "beginner"


class RegisterRequest(BaseModel):
    email: str
    password: str
    nickname: str


class LoginRequest(BaseModel):
    email: str
    password: str


class CompleteLessonRequest(BaseModel):
    session_id: int
    lesson_id: int


class AnswerRecordRequest(BaseModel):
    user_id: int
    session_id: int
    question_id: int
    user_answer: str
    is_correct: bool


class IncorrectAnswerRequest(BaseModel):
    user_id: int
    subject: str = ""
    question: str = ""
    options: list = []
    answer: str = ""
    explanation: str = ""
    user_answer: str = ""
    session_id: int = 0
    question_id: int = 0
    question_type: str = ""
    difficulty: str = ""
    model_answer: str = ""


@app.get("/")
def read_root():
    return {"message": "백엔드 서버 정상 작동 중!"}


def normalize_generation_mode(value: str) -> str:
    normalized = (value or "ai_only").strip()
    aliases = {
        "aiOnly": "ai_only",
        "materialOnly": "material_only",
        "mixed": "mixed",
        "AI": "ai_only",
        "Document": "material_only",
    }
    return aliases.get(normalized, normalized if normalized in {"ai_only", "material_only", "mixed"} else "ai_only")


def normalize_learning_level(value: str) -> str:
    normalized = (value or "beginner").strip()
    aliases = {
        "초급": "beginner",
        "중급": "intermediate",
        "고급": "advanced",
        "beginner": "beginner",
        "intermediate": "intermediate",
        "advanced": "advanced",
    }
    return aliases.get(normalized, "beginner")


def learning_level_label(value: str) -> str:
    return {
        "beginner": "초급",
        "intermediate": "중급",
        "advanced": "고급",
    }.get(normalize_learning_level(value), "초급")


def generation_mode_label(value: str) -> str:
    return {
        "ai_only": "AI 자체 생성",
        "material_only": "업로드한 자료 기반 생성",
        "mixed": "AI + 자료 혼합 생성",
    }.get(normalize_generation_mode(value), "AI 자체 생성")


def level_policy(value: str) -> str:
    level = normalize_learning_level(value)
    if level == "beginner":
        return "핵심 개념 확인 위주입니다. 객관식과 짧은 단답형 비중을 높이고, 정답이 명확한 쉬운 문제를 만드세요."
    if level == "intermediate":
        return "개념 적용 문제를 포함합니다. 객관식, 단답형, 서술형을 섞고, 예시 상황을 보고 개념을 적용하는 문제를 포함하세요."
    return "응용, 분석, 서술형 비중을 높입니다. 과목 성격이 맞으면 코딩형, SQL형, 계산형, 코드 해석형을 포함하고 단순 암기보다 이해와 적용을 요구하세요."


def normalize_curriculum(parsed_curriculum: List[dict]) -> List[dict]:
    if not parsed_curriculum:
        parsed_curriculum = [
            {"level": 1, "title": "기초 개념 확인", "description": "핵심 용어와 전체 구조를 파악합니다."},
            {"level": 2, "title": "주요 개념 이해", "description": "중요 개념과 원리를 정리합니다."},
            {"level": 3, "title": "개념 적용", "description": "예시 상황에 개념을 적용합니다."},
            {"level": 4, "title": "종합 복습", "description": "응용 문제로 전체 내용을 점검합니다."},
        ]

    final_curriculum = []
    for i, item in enumerate(parsed_curriculum):
        if not isinstance(item, dict):
            continue
        try:
            level = int(item.get("level", i + 1))
        except (TypeError, ValueError):
            level = i + 1
        final_curriculum.append(
            {
                "id": i + 1,
                "level": level,
                "title": f"Level {i + 1}. {item.get('title', '학습 단원')}",
                "description": item.get("description", "이 단원의 핵심 개념을 학습합니다."),
                "isLocked": i > 0,
                "isCompleted": False,
            }
        )
    return final_curriculum or normalize_curriculum([])


async def generate_curriculum_from_source_text(source_text: str) -> List[dict]:
    prompt = f"""
    당신은 유능한 교육 전문가입니다. 업로드된 참고 자료의 내용을 학습하기 위해 듀오링고 스타일의 점진적인 4단계 커리큘럼을 생성해 주세요.

    [참고자료 요약 또는 일부 내용]
    {source_text[:6000]}

    반드시 아래 JSON object 형식으로만 응답하세요.
    {{"curriculum": [{{"id": 1, "level": 1, "title": "기초 단어 및 개요", "description": "문서에 소개된 기초적인 주요 용어를 공부합니다."}}]}}

    학습자가 본 문서의 개념을 기초부터 단계별로 학습할 수 있도록 4단계(Level 1 ~ 4) 코스로 작성해 주세요.
    """
    response = await client.chat.completions.create(
        model="openai/gpt-oss-120b",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
        response_format={"type": "json_object"},
    )
    data_object = json.loads(response.choices[0].message.content.strip())
    return normalize_curriculum(data_object.get("curriculum", []))


def decode_text_file(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="ignore")


async def extract_text_from_upload_file(file: UploadFile) -> Tuple[str, str]:
    filename = file.filename or "unknown"
    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail=f"{filename}: 비어 있는 파일입니다.")

    file_ext = os.path.splitext(filename)[1].lower()
    try:
        if file_ext == ".txt":
            extracted_text = decode_text_file(content)
        elif file_ext == ".pdf":
            reader = PdfReader(io.BytesIO(content))
            extracted_text = "\n".join(filter(None, (page.extract_text() for page in reader.pages)))
        elif file_ext == ".pptx":
            prs = Presentation(io.BytesIO(content))
            text_list = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_list.append(shape.text.strip())
            extracted_text = "\n".join(text_list)
        elif file_ext == ".ppt":
            raise HTTPException(
                status_code=400,
                detail=f"{filename}: .ppt는 레거시 PowerPoint 형식이라 서버에서 안정적으로 읽을 수 없습니다. .pptx로 변환한 뒤 다시 업로드해 주세요.",
            )
        else:
            raise HTTPException(status_code=400, detail="지원되지 않는 파일 확장자입니다. (.txt, .pdf, .ppt, .pptx)만 선택할 수 있습니다.")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"{filename}: 텍스트 추출 중 오류가 발생했습니다. {exc}")

    if not extracted_text.strip():
        raise HTTPException(status_code=400, detail=f"{filename}: 문서 내에서 추출할 수 있는 텍스트가 없습니다.")

    print(f"📂 파일 텍스트 추출 성공: {filename} (원본: {len(content)} bytes, 추출: {len(extracted_text)} 자)")
    return filename, extracted_text.strip()


async def build_combined_source_text(files: List[UploadFile]) -> Tuple[str, List[str], int]:
    if not files:
        raise HTTPException(status_code=400, detail="업로드할 파일을 1개 이상 선택해 주세요.")

    uploaded_files: List[str] = []
    source_parts: List[str] = []
    for file in files:
        filename, extracted_text = await extract_text_from_upload_file(file)
        uploaded_files.append(filename)
        source_parts.append(f"[파일명: {filename}]\n{extracted_text}")

    combined_source_text = "\n\n".join(source_parts).strip()
    if not combined_source_text:
        raise HTTPException(status_code=400, detail="업로드한 자료에서 추출된 텍스트가 없습니다.")
    return combined_source_text, uploaded_files, len(combined_source_text)


def build_material_subject(uploaded_files: List[str]) -> str:
    return f"[자료] {uploaded_files[0]}" if len(uploaded_files) == 1 else f"[자료] {uploaded_files[0]} 외 {len(uploaded_files) - 1}개"


def build_type_counts(weights: Dict[str, int], total_count: int) -> Dict[str, int]:
    total_count = max(1, min(total_count, 30))
    if not weights:
        return {"multiple_choice": total_count}

    positive = {key: value for key, value in weights.items() if value > 0}
    total_weight = sum(positive.values())
    if total_weight <= 0:
        return {"multiple_choice": total_count}

    raw = []
    for q_type, weight in positive.items():
        exact = total_count * weight / total_weight
        floor_value = int(exact)
        raw.append({"type": q_type, "count": floor_value, "remain": exact - floor_value})

    used = sum(item["count"] for item in raw)
    remain_count = total_count - used
    raw.sort(key=lambda item: item["remain"], reverse=True)

    for i in range(remain_count):
        raw[i % len(raw)]["count"] += 1

    return {item["type"]: item["count"] for item in raw if item["count"] > 0}


def ensure_list(value) -> List[str]:
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()]
    if value is None:
        return []
    text = str(value).strip()
    return [text] if text else []


def first_non_empty(*values) -> str:
    for value in values:
        text = str(value).strip() if value is not None else ""
        if text:
            return text
    return ""


def build_fallback_model_answer(question_text: str, answer: str = "") -> str:
    if answer:
        return answer
    return f"이 문제는 '{question_text}'에 대해 핵심 개념, 근거, 결론을 포함해 설명하는 답안을 요구합니다."


def normalize_short_answer_question(q: dict) -> dict:
    answer = first_non_empty(q.get("answer"), q.get("expected_answer"), q.get("model_answer"))
    if not answer:
        answer = "핵심 개념"

    acceptable_answers = ensure_list(q.get("acceptable_answers") or q.get("acceptableAnswers"))
    acceptable_answers = list(dict.fromkeys([answer, *acceptable_answers]))

    q["answer"] = answer
    q["acceptable_answers"] = acceptable_answers
    q["options"] = []
    q["model_answer"] = ""
    q["sample_answer"] = ""
    q["expected_answer"] = answer
    q["keywords"] = ensure_list(q.get("keywords"))
    q["grading_criteria"] = first_non_empty(
        q.get("grading_criteria"),
        q.get("gradingCriteria"),
        "한 단어 또는 짧은 구문이 허용 답안과 의미상 같으면 정답으로 인정합니다.",
    )
    q["expected_format"] = first_non_empty(q.get("expected_format"), "한 단어 또는 짧은 구문")
    return q


def is_valid_short_answer_question(q: dict) -> bool:
    question_text = str(q.get("question", ""))
    answer = first_non_empty(q.get("answer"), q.get("expected_answer"), q.get("model_answer"))
    if not answer:
        return False
    if len(answer) > 45 or len(answer.split()) > 8:
        return False
    banned_phrases = ["설명하시오", "서술하시오", "비교하시오", "분석하시오", "논하시오", "이유를 쓰", "과정을 설명"]
    return not any(phrase in question_text for phrase in banned_phrases)

def contains_korean(text: str) -> bool:
    return any("\uac00" <= ch <= "\ud7a3" for ch in str(text or ""))


def is_korean_generated_question(q: dict) -> bool:
    """
    AI가 영어 문제를 내는 경우를 막기 위한 최소 검증입니다.
    Java, SQL, API 같은 기술 용어는 허용하지만,
    문제 본문과 해설에는 반드시 한글이 포함되어야 합니다.
    """
    question = str(q.get("question", ""))
    explanation = str(q.get("explanation", ""))

    if not contains_korean(question):
        return False

    if explanation.strip() and not contains_korean(explanation):
        return False

    options = q.get("options", [])
    if isinstance(options, list) and options:
        joined_options = " ".join(str(option) for option in options)
        if not contains_korean(joined_options):
            # 연산자 문제처럼 &&, ||, == 만 있는 선택지는 허용
            operator_like = all(
                str(option).strip() in {"&&", "||", "!", "==", "!=", "<", ">", "<=", ">=", "+", "-", "*", "/", "%"}
                for option in options
            )
            if not operator_like:
                return False

    return True

def normalize_descriptive_question(q: dict) -> dict:
    answer = first_non_empty(
        q.get("model_answer"),
        q.get("modelAnswer"),
        q.get("sample_answer"),
        q.get("sampleAnswer"),
        q.get("expected_answer"),
        q.get("expectedAnswer"),
        q.get("answer"),
    )
    answer = build_fallback_model_answer(q.get("question", "서술형 문제"), answer)

    keywords = ensure_list(q.get("keywords"))
    if not keywords:
        keywords = ensure_list(q.get("rubric"))[:4]

    grading_criteria = first_non_empty(
        q.get("grading_criteria"),
        q.get("gradingCriteria"),
        "핵심 키워드와 의미가 포함되어 있으면 부분 정답 이상으로 인정합니다.",
    )

    q["answer"] = answer
    q["model_answer"] = answer
    q["sample_answer"] = first_non_empty(q.get("sample_answer"), answer)
    q["expected_answer"] = first_non_empty(q.get("expected_answer"), answer)
    q["keywords"] = keywords
    q["grading_criteria"] = grading_criteria
    q["rubric"] = ensure_list(q.get("rubric")) or keywords
    q["options"] = []
    return q


def normalize_generated_question(q: dict, q_type: str, request: QuestionRequest, source_label: str) -> dict:
    if q_type == "short_answer":
        q = normalize_short_answer_question(q)
    elif q_type == "descriptive":
        q = normalize_descriptive_question(q)
    elif q_type != "multiple_choice":
        q["options"] = []
        q["answer"] = first_non_empty(q.get("answer"), q.get("expected_answer"), q.get("model_answer"))
        q["acceptable_answers"] = ensure_list(q.get("acceptable_answers") or q.get("acceptableAnswers"))
        q["keywords"] = ensure_list(q.get("keywords"))
        q["grading_criteria"] = first_non_empty(q.get("grading_criteria"), q.get("gradingCriteria"))

    return {
        "id": 0,
        "type": q_type,
        "question": q.get("question", "문제를 생성하지 못했습니다."),
        "options": q.get("options", []),
        "answer": q.get("answer", ""),
        "acceptable_answers": ensure_list(q.get("acceptable_answers") or q.get("acceptableAnswers")),
        "explanation": q.get("explanation", "해설이 제공되지 않았습니다."),
        "source_type": source_label,
        "difficulty": learning_level_label(request.learning_level or request.difficulty),
        "code": q.get("code", ""),
        "language": q.get("language", ""),
        "starter_code": q.get("starter_code", ""),
        "test_cases": q.get("test_cases", []),
        "rubric": ensure_list(q.get("rubric")),
        "expected_format": q.get("expected_format", ""),
        "model_answer": q.get("model_answer", ""),
        "sample_answer": q.get("sample_answer", ""),
        "expected_answer": q.get("expected_answer", ""),
        "keywords": ensure_list(q.get("keywords")),
        "grading_criteria": q.get("grading_criteria", ""),
    }


def build_question_prompt(request: QuestionRequest, source_text: Optional[str], type_counts: Dict[str, int]) -> str:
    context_str = f"학습 단원은 '{request.level_title}' ({request.level_description}) 입니다." if request.level_title else ""
    generation_mode = normalize_generation_mode(request.generation_mode)
    source_clause = ""

    if generation_mode == "ai_only":
        basis = f"일반적으로 알려진 {request.subject} 과목 지식을 기반으로 문제를 생성하세요. 업로드 자료가 있더라도 사용하지 마세요."
    elif generation_mode == "material_only":
        basis = f"""
        아래 [참고자료] 내용에만 기반해 문제를 생성하세요. 참고자료에 없는 사실을 임의로 확장하지 마세요.
        [참고자료]
        {source_text[:15000] if source_text else ''}
        """
        source_clause = "참고자료의 표현과 핵심 내용을 우선 반영합니다."
    else:
        basis = f"""
        아래 [참고자료]를 우선 사용하되, 필요한 경우 {request.subject}의 일반 지식으로 맥락을 보완하세요.
        [참고자료]
        {source_text[:15000] if source_text else ''}
        """
        source_clause = "자료의 핵심 내용과 AI의 보충 설명을 균형 있게 섞습니다."

    return f"""
    당신은 유능한 교육 전문가입니다. Flutter 앱이 바로 파싱할 수 있도록 반드시 JSON object만 반환하세요.

    [과목]
    {request.subject}

    [과목 유형]
    {request.subject_type}

    [문제 생성 방식]
    {generation_mode_label(generation_mode)}
    {source_clause}

    [학습 수준]
    {learning_level_label(request.learning_level or request.difficulty)}
    {level_policy(request.learning_level or request.difficulty)}

    [문제 유형 모드]
    {request.mode}

    [학습 단원]
    {context_str}

    [문제 생성 근거]
    {basis}

    [허용된 문제 유형]
    {json.dumps(request.allowed_question_types, ensure_ascii=False)}

    [반드시 지킬 문제 유형별 개수]
    {json.dumps(type_counts, ensure_ascii=False)}

    [중요 규칙]
    1. questions 배열에는 총 {request.count}개의 문제를 넣으세요.
    2. 각 문제의 type은 반드시 허용된 문제 유형 중 하나여야 합니다.
    3. conceptual 과목에는 coding, code_reading, sql_writing, command_writing 문제를 만들지 마세요.
    4. calculation 과목에는 coding 대신 calculation 문제를 사용하세요.
    5. practical 과목 중 데이터베이스/SQL 단원은 sql_writing 문제를 사용할 수 있습니다.
    6. programming 과목은 code_reading 또는 coding 문제를 사용할 수 있습니다.
    7. 객관식 문제만 options를 정확히 4개 제공합니다.
    8. 단답형, 서술형, 코딩형, SQL 작성형, 계산형의 options는 빈 배열로 둡니다.
    9. 단답형(short_answer)은 반드시 한 단어 또는 짧은 구문으로 답할 수 있게 만드세요.
       - 빈칸 채우기, 핵심 개념명 답하기, 간단한 용어 답하기, 짧은 계산 결과 답하기 위주로 만드세요.
       - "설명하시오", "서술하시오", "비교하시오", "분석하시오"처럼 긴 문장 답변을 요구하는 단답형 문제는 금지합니다.
       - short_answer에는 answer와 acceptable_answers를 반드시 넣고 expected_format은 "한 단어 또는 짧은 구문"으로 둡니다.
    10. 서술형(descriptive)은 model_answer, sample_answer, expected_answer 중 최소 하나를 반드시 채우고,
        keywords와 grading_criteria를 반드시 제공합니다.
    11. 서술형 문제에는 rubric을 2~4개 제공합니다.
    12. 코딩형/SQL 작성형 문제에는 starter_code 또는 test_cases를 가능한 경우 제공합니다.
    13. 고급 수준에서는 단순 암기보다 분석, 적용, 코드 해석, SQL/계산 문제를 우선합니다.
    14. 반드시 JSON object만 반환하세요. JSON 밖에는 아무 문장도 쓰지 마세요.
    15. 모든 문제(question), 선택지(options), 정답(answer), 허용답안(acceptable_answers), 해설(explanation), 모범답안(model_answer, sample_answer, expected_answer)은 반드시 한국어로 작성하세요.
    16. 영어 문장으로 된 문제 출제는 금지합니다.
    17. Java, SQL, Flutter, HTTP, API, CPU, RAM 같은 기술 용어와 코드, 연산자, 명령어는 원문을 유지할 수 있지만 설명 문장은 반드시 한국어여야 합니다.
    18. question 필드에 한글이 하나도 없으면 잘못된 문제입니다.
    19. 초급 문제에서는 가능한 한 “아닌 것은?”, “틀린 것은?” 같은 부정형 객관식을 줄이고, 정답이 명확한 개념 확인 문제를 우선 생성하세요.

    [반환 JSON 형식]
    {{
      "questions": [
        {{
          "id": 1,
          "type": "short_answer",
          "question": "객체지향 프로그래밍의 4대 특징 중 내부 구현을 숨기고 필요한 기능만 제공하는 개념은?",
          "options": [],
          "answer": "캡슐화",
          "acceptable_answers": ["캡슐화", "encapsulation"],
          "explanation": "캡슐화는 내부 구현을 숨기고 외부에는 필요한 기능만 제공하는 객체지향 개념입니다.",
          "model_answer": "",
          "sample_answer": "",
          "expected_answer": "캡슐화",
          "keywords": [],
          "grading_criteria": "허용 답안 중 하나와 의미가 같으면 정답으로 인정합니다.",
          "code": "",
          "language": "",
          "starter_code": "",
          "test_cases": [],
          "rubric": [],
          "expected_format": "한 단어 또는 짧은 구문"
        }},
        {{
          "id": 2,
          "type": "descriptive",
          "question": "데이터베이스에서 정규화를 사용하는 이유를 설명하시오.",
          "options": [],
          "answer": "정규화는 데이터 중복을 줄이고 삽입, 삭제, 수정 이상을 방지하여 데이터의 일관성과 무결성을 높이기 위해 사용한다.",
          "acceptable_answers": [],
          "explanation": "정규화의 핵심 목적은 중복 감소와 이상 현상 방지입니다.",
          "model_answer": "정규화는 데이터 중복을 줄이고 삽입, 삭제, 수정 이상을 방지하여 데이터의 일관성과 무결성을 높이기 위해 사용한다.",
          "sample_answer": "정규화는 데이터 중복을 줄이고 이상 현상을 방지하기 위해 사용한다.",
          "expected_answer": "데이터 중복 감소와 이상 현상 방지 목적을 설명한다.",
          "keywords": ["데이터 중복", "이상 현상", "무결성", "일관성"],
          "grading_criteria": "데이터 중복 감소와 이상 현상 방지 목적을 설명하면 정답 또는 부분 정답으로 인정한다.",
          "code": "",
          "language": "",
          "starter_code": "",
          "test_cases": [],
          "rubric": ["데이터 중복 감소", "이상 현상 방지", "일관성 또는 무결성 언급"],
          "expected_format": "2~4문장"
        }}
      ]
    }}
    """


def resolve_question_source(request: QuestionRequest) -> Tuple[Optional[str], str]:
    generation_mode = normalize_generation_mode(request.generation_mode)
    if generation_mode == "ai_only":
        return None, "AI"

    source_text = database.get_session_source_text(request.session_id) if request.session_id > 0 else None
    if not source_text:
        raise HTTPException(status_code=400, detail="자료 기반 또는 혼합 생성은 업로드 자료가 연결된 학습 세션이 필요합니다.")

    if generation_mode == "material_only":
        return source_text, "Document"
    return source_text, "AI+Document"


@app.post("/api/auth/register")
async def register(request: RegisterRequest):
    user = database.register_user(request.email, request.password, request.nickname)
    if not user:
        raise HTTPException(status_code=400, detail="이미 가입된 이메일이거나 회원가입에 실패했습니다.")
    print(f"👤 회원가입 성공: {user['email']} ({user['nickname']})")
    return {"user": user}


@app.post("/api/auth/login")
async def login(request: LoginRequest):
    user = database.login_user(request.email, request.password)
    if not user:
        raise HTTPException(status_code=401, detail="이메일 또는 비밀번호가 일치하지 않습니다.")
    print(f"🔓 로그인 성공: {user['email']} ({user['nickname']})")
    return {"user": user}


@app.get("/api/sessions")
async def get_sessions(user_id: int):
    return {"sessions": database.get_study_sessions(user_id)}


@app.patch("/api/sessions/{session_id}/setup")
async def update_session_setup(session_id: int, request: SessionSetupRequest):
    difficulty = normalize_learning_level(request.learning_level or request.difficulty)
    generation_mode = normalize_generation_mode(request.generation_mode)
    success = database.update_session_setup(session_id, generation_mode, difficulty)
    if not success:
        raise HTTPException(status_code=404, detail="학습 세션을 찾을 수 없습니다.")
    saved_session = database.get_session_by_id(session_id)
    return {"session": saved_session}


@app.post("/api/complete-lesson")
async def complete_lesson(request: CompleteLessonRequest):
    success = database.update_lesson_completion(request.session_id, request.lesson_id)
    if not success:
        raise HTTPException(status_code=400, detail="레벨 완료 처리에 실패했습니다.")
    return {"message": "레벨 완료 처리 및 다음 레벨 잠금 해제 성공"}


@app.post("/api/generate-curriculum")
async def generate_curriculum(request: CurriculumRequest):
    try:
        generation_mode = normalize_generation_mode(request.generation_mode)
        difficulty = normalize_learning_level(request.learning_level or request.difficulty)

        existing_session = database.get_session_by_subject(request.user_id, request.subject)
        if existing_session:
            print(f"♻️ 기존 '{request.subject}' 학습 세션을 재사용합니다.")
            database.update_session_setup(existing_session["id"], generation_mode, difficulty)
            saved_session = database.get_session_by_id(existing_session["id"])
            return {
                "session_id": saved_session["id"],
                "subject": saved_session["subject"],
                "progress": saved_session["progress"],
                "generation_mode": saved_session["generation_mode"],
                "difficulty": saved_session["difficulty"],
                "last_studied_at": saved_session.get("last_studied_at"),
                "curriculum": saved_session["curriculum"],
            }

        prompt = f"""
        당신은 유능한 교육 전문가입니다. 사용자가 학습하고자 하는 주제: '{request.subject}' 에 대한 적절성 검사 및 듀오링고 스타일의 점진적인 4단계 커리큘럼을 생성해 주세요.
        입력값이 무의미한 자판 배열, 스팸, 심한 비속어라면 valid=false와 친절한 error_message를 반환하세요.
        학습 수준은 {learning_level_label(difficulty)}입니다. {level_policy(difficulty)}
        반드시 JSON object 형식으로만 응답하세요.
        {{"valid": true, "error_message": "", "curriculum": [{{"id": 1, "level": 1, "title": "기초 다지기", "description": "이 주제에 대한 기본 용어와 개념을 이해합니다."}}]}}
        학습자가 기초부터 응용까지 올라갈 수 있도록 4단계(Level 1 ~ 4)로 구성해 주세요.
        """
        response = await client.chat.completions.create(
            model="openai/gpt-oss-120b",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
            response_format={"type": "json_object"},
        )
        data_object = json.loads(response.choices[0].message.content.strip())
        if not data_object.get("valid", True):
            error_msg = data_object.get("error_message", "올바른 학습 주제가 아닙니다. 배우고 싶은 과목이나 주제를 명확하게 적어주세요!")
            print(f"⚠️ 학습 주제 무효 처리: '{request.subject}' -> {error_msg}")
            raise HTTPException(status_code=400, detail=error_msg)

        session_id = database.create_study_session(
            request.user_id,
            request.subject,
            normalize_curriculum(data_object.get("curriculum", [])),
            None,
            generation_mode,
            difficulty,
        )
        saved_session = database.get_session_by_id(session_id)
        print(f"✅ 성공적으로 '{request.subject}'에 대한 커리큘럼 세션(ID: {session_id})을 생성 및 저장했습니다.")
        return {
            "session_id": saved_session["id"],
            "subject": saved_session["subject"],
            "progress": saved_session["progress"],
            "generation_mode": saved_session["generation_mode"],
            "difficulty": saved_session["difficulty"],
            "last_studied_at": saved_session.get("last_studied_at"),
            "curriculum": saved_session["curriculum"],
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ 커리큘럼 생성 중 에러 발생: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/generate-questions")
async def generate_questions(request: QuestionRequest):
    try:
        request.generation_mode = normalize_generation_mode(request.generation_mode)
        request.learning_level = normalize_learning_level(request.learning_level or request.difficulty)
        request.difficulty = request.learning_level

        source_text, source_label = resolve_question_source(request)
        if source_text:
            print(f"📖 세션 ID {request.session_id}의 문서 데이터를 활용하여 문제를 생성합니다. ({source_label})")
        else:
            print(f"🤖 일반 지식 기반 문제를 생성합니다: '{request.subject}'")

        if not request.allowed_question_types:
            request.allowed_question_types = ["multiple_choice"]
        if not request.question_type_weights:
            request.question_type_weights = {"multiple_choice": 100}

        allowed_set = set(request.allowed_question_types)
        type_counts = build_type_counts(request.question_type_weights, request.count)
        prompt = build_question_prompt(request, source_text, type_counts)

        response = await client.chat.completions.create(
            model="openai/gpt-oss-120b",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2,
            response_format={"type": "json_object"},
        )
        parsed_questions = json.loads(response.choices[0].message.content.strip()).get("questions", [])
        final_questions = []
        for q in parsed_questions:
            if len(final_questions) >= request.count:
                break
            if isinstance(q, str):
                try:
                    q = json.loads(q)
                except json.JSONDecodeError:
                    continue
            if not isinstance(q, dict):
                continue

            q_type = q.get("type", "multiple_choice")
            if q_type not in allowed_set:
                continue

            options = q.get("options", [])
            if q_type == "multiple_choice":
                if not isinstance(options, list) or len(options) != 4:
                    continue
            elif q_type == "short_answer" and not is_valid_short_answer_question(q):
                continue
            else:
                q["options"] = []

            if not is_korean_generated_question(q):
                continue

            normalized_question = normalize_generated_question(q, q_type, request, source_label)
            normalized_question["id"] = len(final_questions) + 1
            final_questions.append(normalized_question)

        actual_count = len(final_questions)
        if actual_count < request.count:
            print(f"⚠️ 요청한 {request.count}문제보다 적은 {actual_count}문제만 생성되었습니다.")
        else:
            print(f"✅ 성공적으로 {actual_count}문제를 생성하여 반환합니다.")

        saved_questions = database.save_generated_questions(request.session_id, final_questions)
        return {"questions": saved_questions}
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ API 요청 처리 중 내부 에러 발생: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/answer-records")
async def add_answer_record(request: AnswerRecordRequest):
    row_id = database.add_answer_record(
        request.user_id,
        request.session_id,
        request.question_id,
        request.user_answer,
        request.is_correct,
    )
    if not row_id:
        raise HTTPException(status_code=400, detail="풀이 기록을 저장하지 못했습니다.")
    return {"id": row_id, "message": "풀이 기록 저장 완료"}


@app.get("/api/incorrect-answers")
async def get_incorrect_answers(user_id: int):
    return {"answers": database.get_incorrect_answers(user_id)}


@app.post("/api/incorrect-answers")
async def add_incorrect_answer(request: IncorrectAnswerRequest):
    row_id = database.add_incorrect_answer(
        request.user_id,
        request.subject,
        request.question,
        request.options,
        request.answer,
        request.explanation,
        request.user_answer,
        request.session_id,
        request.question_id,
        request.question_type,
        request.difficulty,
        request.model_answer,
    )
    if not row_id:
        raise HTTPException(status_code=500, detail="오답을 DB에 등록하지 못했습니다.")
    print(f"📝 오답 추가 완료: (User ID: {request.user_id}, 과목: {request.subject}, Question ID: {request.question_id})")
    return {"id": row_id, "message": "오답 노트 등록 완료"}


@app.delete("/api/incorrect-answers/{answer_id}")
async def delete_incorrect_answer(answer_id: int):
    success = database.delete_incorrect_answer(answer_id)
    if not success:
        raise HTTPException(status_code=404, detail="해당 오답 항목을 찾을 수 없습니다.")
    print(f"🗑️ 오답 복습 완료 처리 (ID: {answer_id})")
    return {"message": "오답 노트에서 복습 완료 처리되었습니다."}

@app.delete("/api/sessions/{session_id}")
async def delete_study_session(session_id: int, user_id: int):
    success = database.delete_study_session(user_id, session_id)
    if not success:
        raise HTTPException(status_code=404, detail="삭제할 학습 세션을 찾을 수 없습니다.")

    print(f"🗑️ 학습 세션 삭제 완료: session_id={session_id}, user_id={user_id}")
    return {"message": "학습 세션이 삭제되었습니다."}


@app.delete("/api/incorrect-answers")
async def clear_incorrect_answers(user_id: int, session_id: Optional[int] = None):
    success = database.clear_incorrect_answers(user_id, session_id)
    if not success:
        raise HTTPException(status_code=500, detail="오답노트 삭제에 실패했습니다.")

    if session_id:
        print(f"🧹 특정 학습 오답노트 삭제 완료: user_id={user_id}, session_id={session_id}")
    else:
        print(f"🧹 전체 오답노트 삭제 완료: user_id={user_id}")

    return {"message": "오답노트가 삭제되었습니다."}

@app.get("/api/study-calendar")
async def get_study_calendar(user_id: int):
    return database.get_study_calendar(user_id)


@app.post("/api/upload-material")
async def upload_material(
    user_id: int = Form(...),
    files: Optional[List[UploadFile]] = File(None),
    file: Optional[UploadFile] = File(None),
    session_id: Optional[int] = Form(None),
    regenerate_curriculum: bool = Form(False),
    generation_mode: str = Form("material_only"),
    difficulty: str = Form("beginner"),
):
    try:
        generation_mode = normalize_generation_mode(generation_mode)
        difficulty = normalize_learning_level(difficulty)

        upload_files: List[UploadFile] = []
        if files:
            upload_files.extend(files)
        if file:
            upload_files.append(file)

        combined_source_text, uploaded_files, total_extracted_length = await build_combined_source_text(upload_files)

        if session_id:
            existing_session = database.get_session_by_id(session_id)
            if not existing_session:
                raise HTTPException(status_code=404, detail="자료를 추가할 기존 로드맵을 찾을 수 없습니다.")
            if existing_session["user_id"] != user_id:
                raise HTTPException(status_code=403, detail="다른 사용자의 로드맵에는 자료를 추가할 수 없습니다.")

            appended = database.append_session_source_text(session_id, combined_source_text)
            if not appended:
                raise HTTPException(status_code=500, detail="기존 로드맵에 자료를 추가하지 못했습니다.")
            database.update_session_setup(session_id, generation_mode, difficulty)
            if regenerate_curriculum:
                updated_source_text = database.get_session_source_text(session_id) or combined_source_text
                regenerated_curriculum = await generate_curriculum_from_source_text(updated_source_text)
                database.replace_session_curriculum(session_id, regenerated_curriculum)

            saved_session = database.get_session_by_id(session_id)
            print(f"✅ 기존 로드맵(ID: {session_id})에 자료 {len(uploaded_files)}개 추가 완료 (커리큘럼 재생성: {regenerate_curriculum})")
            return {
                "session_id": saved_session["id"],
                "subject": saved_session["subject"],
                "progress": saved_session["progress"],
                "generation_mode": saved_session["generation_mode"],
                "difficulty": saved_session["difficulty"],
                "last_studied_at": saved_session.get("last_studied_at"),
                "curriculum": saved_session["curriculum"],
                "uploaded_files": uploaded_files,
                "total_extracted_length": total_extracted_length,
                "appended_to_session": True,
                "regenerated_curriculum": regenerate_curriculum,
            }

        subject_name = build_material_subject(uploaded_files)
        existing_session = database.get_session_by_subject(user_id, subject_name)
        if existing_session:
            print(f"♻️ 기존 업로드 세션 '{subject_name}'을 반환합니다.")
            database.update_session_setup(existing_session["id"], generation_mode, difficulty)
            saved_session = database.get_session_by_id(existing_session["id"])
            return {
                "session_id": saved_session["id"],
                "subject": saved_session["subject"],
                "progress": saved_session["progress"],
                "generation_mode": saved_session["generation_mode"],
                "difficulty": saved_session["difficulty"],
                "last_studied_at": saved_session.get("last_studied_at"),
                "curriculum": saved_session["curriculum"],
                "uploaded_files": uploaded_files,
                "total_extracted_length": total_extracted_length,
                "appended_to_session": False,
                "regenerated_curriculum": False,
            }

        new_session_id = database.create_study_session(
            user_id,
            subject_name,
            await generate_curriculum_from_source_text(combined_source_text),
            combined_source_text,
            generation_mode,
            difficulty,
        )
        saved_session = database.get_session_by_id(new_session_id)
        print(f"✅ 업로드 자료 기반 커리큘럼(ID: {new_session_id}) 생성 성공!")
        return {
            "session_id": saved_session["id"],
            "subject": saved_session["subject"],
            "progress": saved_session["progress"],
            "generation_mode": saved_session["generation_mode"],
            "difficulty": saved_session["difficulty"],
            "last_studied_at": saved_session.get("last_studied_at"),
            "curriculum": saved_session["curriculum"],
            "uploaded_files": uploaded_files,
            "total_extracted_length": total_extracted_length,
            "appended_to_session": False,
            "regenerated_curriculum": False,
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ 자료 업로드/분석 실패: {e}")
        raise HTTPException(status_code=500, detail=str(e))

